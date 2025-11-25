import os
import json
import requests
from pptx import Presentation
from pptx.util import Pt
from dotenv import load_dotenv
load_dotenv()

# -------------------------------
# 1. CONFIGURE GEMINI API (REST)
# -------------------------------
# Make sure you set your Gemini API key in environment variable:
# export GOOGLE_API_KEY="your_key_here"

API_KEY = os.getenv("GOOGLE_API_KEY")
if not API_KEY:
    raise Exception("Set GOOGLE_API_KEY environment variable")

# Using one of the models you listed: models/gemini-2.5-flash
GEMINI_MODEL = "gemini-2.5-flash"
GEMINI_URL = f"https://generativelanguage.googleapis.com/v1beta/models/{GEMINI_MODEL}:generateContent"


# -------------------------------
# 2. SIMPLE GEMINI FUNCTION (REST)
# -------------------------------
def generate_slide(topic: str):
    prompt = f"""
    Create a single presentation slide for the topic: {topic}

    Return JSON ONLY in this exact format:
    {{
      "title": "Slide Title",
      "bullets": ["point 1", "point 2", "point 3"]
    }}
    """

    payload = {
        "contents": [
            {
                "parts": [
                    {"text": prompt}
                ]
            }
        ]
    }

    # Call Gemini via REST
    response = requests.post(
        GEMINI_URL,
        params={"key": API_KEY},
        json=payload,
        timeout=60
    )

    if response.status_code != 200:
        raise Exception(f"Gemini API error: {response.status_code} - {response.text}")

    data_json = response.json()

    # Extract text from first candidate
    try:
        text = data_json["candidates"][0]["content"]["parts"][0]["text"]
    except Exception:
        raise Exception(f"Unexpected response format: {json.dumps(data_json, indent=2)}")

    # Try to parse JSON from the model's text
    try:
        slide_data = json.loads(text)
    except json.JSONDecodeError:
        # fallback: very simple parsing
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        slide_data = {
            "title": lines[0] if lines else "Untitled Slide",
            "bullets": lines[1:4] if len(lines) > 1 else []
        }

    # Ensure bullets is a list
    if not isinstance(slide_data.get("bullets"), list):
        slide_data["bullets"] = [str(slide_data["bullets"])]

    return slide_data


# -------------------------------
# 3. BUILD PPT USING python-pptx
# -------------------------------
def build_pptx(slide_data, ppt_name="output.pptx"):
    prs = Presentation()

    # Title + content layout (usually index 1)
    layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = slide_data.get("title", "Untitled Slide")

    body = slide.placeholders[1].text_frame
    body.clear()

    bullets = slide_data.get("bullets", [])
    if not bullets:
        bullets = ["No bullet points generated."]

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = str(bullet)
        p.font.size = Pt(16)

    prs.save(ppt_name)
    print(f"✅ PPT generated: {ppt_name}")


# -------------------------------
# 4. MAIN TEST FLOW
# -------------------------------
if __name__ == "__main__":
    topic = "Impact of Artificial Intelligence in Healthcare"

    print("🔍 Generating slide using Gemini (REST)...")
    slide_data = generate_slide(topic)

    print("📄 Slide data received:")
    print(slide_data)

    print("🎨 Building PPT using python-pptx...")
    build_pptx(slide_data)
