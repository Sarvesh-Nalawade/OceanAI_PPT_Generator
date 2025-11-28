import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.dml.color import RGBColor

def create_transportation_presentation():
    """
    Generates a PowerPoint presentation about the improvement of human transportation
    in the past 100 years using the python-pptx library.
    """
    # --- Create Presentation ---
    prs = Presentation()
    # Use a 16:9 aspect ratio
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "The Revolution in Motion"
    subtitle.text = "Human Transportation Over the Last 100 Years"

    # --- Slide 2: Introduction ---
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "From Horsepower to Hypersonic"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "The 20th and 21st centuries witnessed an unprecedented leap in mobility."
    
    p = tf.add_paragraph()
    p.text = "This presentation explores the key advancements in:"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Land: The rise of the automobile and high-speed rail."
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "Sea: From ocean liners to massive cargo ships."
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "Air: The dawn of the jet age and accessible global travel."
    p.level = 2

    # --- Slide 3: The Age of the Automobile ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "The Age of the Automobile"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "The car transformed from a luxury item to a daily necessity."
    
    p = tf.add_paragraph()
    p.text = "Mass Production: Henry Ford's Model T made cars affordable."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Infrastructure: Development of extensive highway systems."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Modern Era: Focus on safety, efficiency, and electric power."
    p.level = 1

    # Image Placeholder for Slide 3
    left = Inches(8)
    top = Inches(4)
    width = Inches(7)
    height = Inches(3.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.text = "Image Placeholder:\nA 1920s Ford Model T next to a modern electric car."
    shape.text_frame.paragraphs[0].font.size = Pt(14)
    shape.text_frame.paragraphs[0].font.italic = True
    
    # --- Slide 4: Revolution on Rails ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Revolution on Rails"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Rail travel evolved from steam-powered workhorses to sleek, high-speed networks."
    
    p = tf.add_paragraph()
    p.text = "Steam to Diesel-Electric: Increased power and efficiency in the mid-20th century."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "High-Speed Rail: Japan's Shinkansen (1964) pioneered travel over 200 km/h."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Magnetic Levitation (Maglev): Pushing speeds beyond 600 km/h."
    p.level = 1

    # Image Placeholder for Slide 4
    left = Inches(8)
    top = Inches(4)
    width = Inches(7)
    height = Inches(3.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.text = "Image Placeholder:\nA classic steam locomotive contrasted with a modern Maglev train."
    shape.text_frame.paragraphs[0].font.size = Pt(14)
    shape.text_frame.paragraphs[0].font.italic = True
    
    # --- Slide 5: Table of Speed ---
    slide_layout = prs.slide_layouts[5]  # Title Only layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "A Century of Speed: A Comparison"
    
    rows, cols = 5, 4
    left = Inches(1.5)
    top = Inches(2.0)
    width = Inches(13.0)
    height = Inches(4.0)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Set column widths
    table.columns[0].width = Inches(3.0)
    table.columns[1].width = Inches(3.5)
    table.columns[2].width = Inches(3.5)
    table.columns[3].width = Inches(3.0)
    
    # Write headers
    headers = ['Mode of Transport', '1920s Top Speed (approx.)', '2020s Top Speed (approx.)', 'NY to LA Travel Time']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        
    # Write data
    data = [
        ["Automobile", "100 km/h", "400+ km/h (Hypercar)", "~4-5 days"],
        ["Train", "160 km/h", "600+ km/h (Maglev)", "~2-3 days"],
        ["Ship", "55 km/h", "65 km/h (Cruise Ship)", "~2 weeks (by sea)"],
        ["Airplane", "240 km/h", "950 km/h (Airliner)", "~5-6 hours"]
    ]
    for r, row_data in enumerate(data):
        for c, cell_data in enumerate(row_data):
            table.cell(r + 1, c).text = cell_data

    # --- Slide 6: Conquering the Skies ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Conquering the Skies"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Aviation matured from a dangerous experiment to the backbone of global travel."
    p = tf.add_paragraph()
    p.text = "The Golden Age (1920s-30s): Rapid innovation in aircraft design."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "The Jet Age (1950s): Commercial jets cut travel times in half."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Mass Air Travel: Deregulation and larger aircraft made flying accessible."
    p.level = 1
    
    # --- Slide 7: Bar Chart - Growth of Air Travel ---
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Explosive Growth in Global Air Travel"
    
    chart_data = CategoryChartData()
    chart_data.categories = ['1950', '1970', '1990', '2010', '2019']
    chart_data.add_series('Passengers (in millions)', (31, 383, 1270, 2710, 4540))
    
    x, y, cx, cy = Inches(2), Inches(2), Inches(12), Inches(5.5)
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )
    chart = graphic_frame.chart
    
    chart.has_legend = False
    
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(12)

    value_axis = chart.value_axis
    value_axis.has_major_gridlines = True
    value_axis.tick_labels.font.size = Pt(12)
    
    value_axis.has_title = True
    value_axis.axis_title.text_frame.text = "Passengers (in millions)"

    # --- Slide 8: Pie Chart - Modal Share of Transportation ---
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "How We Travel Today: Passenger Modal Share (Typical)"
    
    chart_data = ChartData()
    chart_data.categories = ['Personal Vehicle', 'Air Travel', 'Rail', 'Bus/Coach', 'Other']
    chart_data.add_series('Modal Share', (0.85, 0.08, 0.02, 0.03, 0.02))
    
    x, y, cx, cy = Inches(4), Inches(1.5), Inches(8), Inches(6)
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    )
    chart = graphic_frame.chart
    
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0%'
    data_labels.position = XL_LABEL_POSITION.BEST_FIT
    data_labels.font.size = Pt(14)
    data_labels.font.bold = True
    
    # --- Slide 9: The Final Frontier: Space Travel ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "The Final Frontier: Space Travel"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "The ultimate expression of transportation advancement."
    p = tf.add_paragraph()
    p.text = "The Space Race: Apollo missions proved interplanetary travel was possible."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Reusable Rockets: Companies like SpaceX have dramatically reduced launch costs."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Commercialization: The dawn of space tourism and private space stations."
    p.level = 1

    # Image Placeholder for Slide 9
    left = Inches(8)
    top = Inches(4)
    width = Inches(7)
    height = Inches(3.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.text = "Image Placeholder:\nThe Saturn V rocket next to a reusable SpaceX Falcon 9."
    shape.text_frame.paragraphs[0].font.size = Pt(14)
    shape.text_frame.paragraphs[0].font.italic = True
    
    # --- Slide 10: The Future of Transportation ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "The Future of Transportation"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Innovation continues to accelerate, promising an even more connected future."
    p = tf.add_paragraph()
    p.text = "Autonomous Vehicles: Self-driving cars, trucks, and drones."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Hyperloop: High-speed travel in low-pressure tubes."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Urban Air Mobility (UAM): Electric Vertical Take-Off and Landing (eVTOL) aircraft."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Sustainable Fuels: Hydrogen and advanced biofuels for planes and ships."
    p.level = 1
    
    # --- Slide 11: Conclusion ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Conclusion"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Over the past century, transportation has evolved at an astonishing pace."
    p = tf.add_paragraph()
    p.text = "Speed, safety, and accessibility have improved beyond imagination."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "The journey continues, driven by technology and the human desire to explore."
    p.level = 1
    
    # --- Slide 12: Thank You ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Thank You"
    subtitle.text = "Questions?"

    # --- Save Presentation ---
    output_filename = 'output.pptx'
    prs.save(output_filename)
    print(f"Presentation '{output_filename}' created successfully.")

if __name__ == '__main__':
    create_transportation_presentation()